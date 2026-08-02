"""PPTX extraction adapter, wrapping ``python-pptx`` (extra ``viparse[office]``).

Slides are walked in order and each shape's runs record their **font name**, which is
how the normalizer recognizes a legacy encoding (``.VnTime`` → TCVN3). The engine
applies no Vietnamese logic itself.

Two things about a presentation that a document does not have:

**Shapes can be grouped, and a group can contain a group.** A walk that only looks at
top-level shapes silently drops everything inside one — the same shape of defect that
lost tracked insertions in DOCX (VIP-95) and footnotes in the notes part (VIP-98). This
recurses.

**Speaker notes are text nobody sees on the slide.** They are included, after the
slide's own content, on the same reasoning as footnotes: text that is in the file and
silently dropped is the worst outcome, and inlining it among the slide body would put
the presenter's asides in the middle of the audience's text.

``python-pptx`` is imported lazily inside :meth:`PptxEngine.extract`, so importing this
module never requires the dependency; only extraction does.
"""

from __future__ import annotations

from collections.abc import Iterator
from typing import Any

from viparse.detect import CONTENT_TYPE_PPTX
from viparse.engines._shared import blocks_to_text
from viparse.errors import MissingDependency
from viparse.model import RawExtraction
from viparse.options import LoadOptions
from viparse.protocols import DEFAULT_PRIORITY, Source

_INSTALL_HINT = (
    "python-pptx is required for PPTX extraction; install it with: pip install 'viparse[office]'"
)


def _import_pptx() -> Any:
    """Import ``python-pptx`` lazily, raising a clear error if it is missing."""
    try:
        import pptx
    except ImportError as exc:
        raise MissingDependency(_INSTALL_HINT) from exc
    return pptx


class PptxEngine:
    """Extracts ordered text, run fonts, and block structure from a ``.pptx`` file."""

    priority = DEFAULT_PRIORITY
    dependency = "pptx"
    extra = "office"

    def supports(self, content_type: str) -> bool:
        return content_type == CONTENT_TYPE_PPTX

    def extract(self, source: Source, options: LoadOptions) -> RawExtraction:
        pptx = _import_pptx()
        presentation = pptx.Presentation(str(source))
        blocks: list[dict[str, Any]] = []
        fonts: set[str] = set()

        for index, slide in enumerate(presentation.slides, start=1):
            for shape in _iter_shapes(slide.shapes):
                blocks.extend(_shape_blocks(shape, slide, fonts))
            blocks.extend(_notes_blocks(slide, index, fonts))

        return RawExtraction(
            source=str(source),
            content_type=CONTENT_TYPE_PPTX,
            text=blocks_to_text(blocks),
            engine="pptx",
            signals={"fonts": sorted(fonts), "blocks": blocks},
        )


def _iter_shapes(shapes: Any) -> Iterator[Any]:
    """Every shape in document order, descending into groups.

    A group can nest, and PowerPoint files from templates routinely put whole content
    areas inside one. Yielding only the top level would drop them without a word.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _run_font(run: Any, paragraph: Any) -> str | None:
    """The run's font, falling back to the paragraph's.

    A run that inherits its typeface reports ``None``, and in a legacy presentation the
    legacy font is often set once on the paragraph rather than on each run — so reading
    the run alone loses the signal the normalizer needs.
    """
    name = run.font.name
    if name:
        return str(name)
    paragraph_font = getattr(paragraph, "font", None)
    inherited = getattr(paragraph_font, "name", None) if paragraph_font is not None else None
    return str(inherited) if inherited else None


def _paragraph_block(paragraph: Any, is_title: bool) -> dict[str, Any] | None:
    """Map one paragraph to a heading/paragraph block, or ``None`` if it is empty."""
    runs = [
        {"text": run.text, "font": _run_font(run, paragraph)} for run in paragraph.runs if run.text
    ]
    text = "".join(str(run["text"]) for run in runs)
    if not text.strip():
        return None
    if is_title:
        block: dict[str, Any] = {"type": "heading", "level": 1, "text": text}
    else:
        # Outline level maps to nesting, and a nested bullet is still a paragraph — the
        # renderer keeps structure, not slide geometry.
        block = {"type": "paragraph", "text": text}
    if runs:
        block["runs"] = runs
    return block


def _shape_blocks(shape: Any, slide: Any, fonts: set[str]) -> list[dict[str, Any]]:
    """Blocks for one shape: a table, a text frame, or nothing."""
    blocks: list[dict[str, Any]] = []

    if getattr(shape, "has_table", False):
        rows: list[list[str]] = []
        table_fonts: set[str] = set()
        for row in shape.table.rows:
            cells = []
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        name = _run_font(run, paragraph)
                        if name:
                            table_fonts.add(name)
                cells.append(cell.text)
            rows.append(cells)
        fonts.update(table_fonts)
        table: dict[str, Any] = {"type": "table", "rows": rows}
        if table_fonts:
            table["fonts"] = sorted(table_fonts)
        return [table]

    if not getattr(shape, "has_text_frame", False):
        return blocks

    title = getattr(slide.shapes, "title", None)
    is_title = title is not None and shape is title
    for paragraph in shape.text_frame.paragraphs:
        block = _paragraph_block(paragraph, is_title)
        if block is None:
            continue
        block_fonts = {str(run["font"]) for run in block.get("runs", []) if run.get("font")}
        if block_fonts:
            block["fonts"] = sorted(block_fonts)
            fonts.update(block_fonts)
        blocks.append(block)
        is_title = False  # only the first paragraph of a title shape is the heading
    return blocks


def _notes_blocks(slide: Any, index: int, fonts: set[str]) -> list[dict[str, Any]]:
    """Speaker notes, after the slide they belong to and labelled with its number.

    Labelled because notes and slide text read identically once flattened, and a reader
    who cannot tell them apart is worse served than one who never saw the notes.
    """
    if not slide.has_notes_slide:
        return []
    frame = slide.notes_slide.notes_text_frame
    if frame is None or not frame.text.strip():
        return []
    blocks: list[dict[str, Any]] = []
    for paragraph in frame.paragraphs:
        block = _paragraph_block(paragraph, is_title=False)
        if block is None:
            continue
        block_fonts = {str(run["font"]) for run in block.get("runs", []) if run.get("font")}
        if block_fonts:
            block["fonts"] = sorted(block_fonts)
            fonts.update(block_fonts)
        blocks.append(block)
    if blocks:
        blocks.insert(0, {"type": "heading", "level": 2, "text": f"[Notes, slide {index}]"})
    return blocks
