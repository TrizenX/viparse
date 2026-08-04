"""PPTX extraction (VIP-109).

A presentation has two structures a document does not: shapes can be **grouped**, and a
slide can carry **speaker notes** nobody sees on screen. Both are places text goes
missing quietly, which is the defect this codebase has now fixed three times in other
containers — so both are tested here rather than assumed.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

import pytest

from viparse.detect import CONTENT_TYPE_PPTX
from viparse.engines.pptx import PptxEngine
from viparse.options import LoadOptions

pytest.importorskip("pptx")

_TCVN3_TITLE = "B¸o c¸o tµi chÝnh"
_TCVN3_BODY = "QuyÕt ®Þnh cña Bé tr­ëng"


def _presentation() -> Any:
    import pptx

    return pptx.Presentation()


def _blank_slide(presentation: Any) -> Any:
    return presentation.slides.add_slide(presentation.slide_layouts[6])


def _textbox(slide: Any, text: str, font: str | None = None) -> Any:
    from pptx.util import Inches

    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    frame = shape.text_frame
    frame.text = text
    if font:
        frame.paragraphs[0].runs[0].font.name = font
    return shape


def _extract(presentation: Any, tmp_path: Path, name: str) -> Any:
    source = tmp_path / name
    presentation.save(str(source))
    return PptxEngine().extract(source, LoadOptions())


def test_supports_only_pptx() -> None:
    engine = PptxEngine()
    assert engine.supports(CONTENT_TYPE_PPTX)
    assert not engine.supports("application/pdf")


def test_slide_text_and_font_signal_are_extracted(tmp_path: Path) -> None:
    """The font name is the whole point: it is how a legacy encoding is recognised."""
    presentation = _presentation()
    _textbox(_blank_slide(presentation), _TCVN3_BODY, ".VnTime")
    raw = _extract(presentation, tmp_path, "basic.pptx")

    assert _TCVN3_BODY in raw.text
    assert ".VnTime" in raw.signals["fonts"]


def test_text_inside_a_grouped_shape_is_extracted(tmp_path: Path) -> None:
    """A walk over top-level shapes only would drop this silently.

    The same shape of defect lost tracked insertions in DOCX and footnotes in their own
    OOXML part. Templates routinely put whole content areas inside a group.
    """
    presentation = _presentation()
    slide = _blank_slide(presentation)
    first = _textbox(slide, "outside the group")
    second = _textbox(slide, _TCVN3_BODY, ".VnTime")
    slide.shapes.add_group_shape([first, second])
    raw = _extract(presentation, tmp_path, "grouped.pptx")

    assert _TCVN3_BODY in raw.text
    assert ".VnTime" in raw.signals["fonts"]


def test_table_cells_are_extracted(tmp_path: Path) -> None:
    from pptx.util import Inches

    presentation = _presentation()
    slide = _blank_slide(presentation)
    table = slide.shapes.add_table(1, 2, Inches(1), Inches(1), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "ChØ tiªu"
    table.cell(0, 1).text = "Sè liÖu"
    raw = _extract(presentation, tmp_path, "table.pptx")

    assert "ChØ tiªu" in raw.text
    assert "Sè liÖu" in raw.text


def test_speaker_notes_are_extracted_and_labelled(tmp_path: Path) -> None:
    """Notes are text in the file; dropping them silently is the worst outcome.

    Labelled and placed after the slide, because notes and slide text read identically
    once flattened — a reader who cannot tell them apart is worse served than one who
    never saw them.
    """
    presentation = _presentation()
    slide = _blank_slide(presentation)
    _textbox(slide, "slide body")
    slide.notes_slide.notes_text_frame.text = "presenter aside"
    raw = _extract(presentation, tmp_path, "notes.pptx")

    assert "presenter aside" in raw.text
    assert "[Notes, slide 1]" in raw.text
    assert raw.text.index("slide body") < raw.text.index("presenter aside")


def test_slides_keep_their_order(tmp_path: Path) -> None:
    presentation = _presentation()
    for text in ("first", "second", "third"):
        _textbox(_blank_slide(presentation), text)
    raw = _extract(presentation, tmp_path, "order.pptx")

    assert raw.text.index("first") < raw.text.index("second") < raw.text.index("third")


def test_a_presentation_with_no_notes_gains_no_notes_block(tmp_path: Path) -> None:
    presentation = _presentation()
    _textbox(_blank_slide(presentation), "just a slide")
    raw = _extract(presentation, tmp_path, "nonotes.pptx")

    assert "[Notes" not in raw.text


def test_legacy_pptx_converts_end_to_end(tmp_path: Path) -> None:
    """The engine reports a font; the normalizer does the rest."""
    import viparse

    presentation = _presentation()
    slide = _blank_slide(presentation)
    _textbox(slide, _TCVN3_TITLE, ".VnTimeH")
    source = tmp_path / "legacy.pptx"
    presentation.save(str(source))

    document = viparse.load(str(source), output="text")[0]
    assert document.metadata.encoding_detected == "tcvn3"
    assert "Báo cáo tài chính" in document.text


def test_slide_title_becomes_a_heading(tmp_path: Path) -> None:
    """The bug the structure benchmark found: titles came out as ordinary paragraphs.

    ``shape is slide.shapes.title`` never matched, because python-pptx builds a fresh
    proxy on every access — ``slide.shapes.title is slide.shapes.title`` is itself
    ``False``. The title was always present in the text, just never marked, so nothing
    downstream had a section to chunk on.
    """
    presentation = _presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Tình hình kinh tế vĩ mô"
    slide.placeholders[1].text_frame.text = "Đoạn số 01."

    blocks = _extract(presentation, tmp_path, "titled.pptx").signals["blocks"]
    headings = [b for b in blocks if b["type"] == "heading"]
    assert [h["text"] for h in headings] == ["Tình hình kinh tế vĩ mô"]
    assert headings[0]["level"] == 1
    # The body is still a paragraph, not swept into the heading.
    assert any(b["type"] == "paragraph" and b["text"] == "Đoạn số 01." for b in blocks)


def test_only_the_title_shape_becomes_a_heading(tmp_path: Path) -> None:
    # A second text shape on a titled slide must not be promoted just by being present.
    presentation = _presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Số liệu"
    _textbox(slide, "Ghi chú bên lề")

    blocks = _extract(presentation, tmp_path, "one_title.pptx").signals["blocks"]
    assert [b["text"] for b in blocks if b["type"] == "heading"] == ["Số liệu"]
